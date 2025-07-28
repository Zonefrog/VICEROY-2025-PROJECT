from pptx import Presentation
from slide_class import SlideData
from logging_funcs import print_, print_2
from pptx.util import Pt

#configurable variables
MAX_SOURCE_IDS_PER_SLIDE = 6

DEFAULT_FONT_SIZE_PT = 24
MIN_FONT_SIZE_PT = 12
FONT_DECREMENT_PT = 1
MAX_LINES_PER_SLIDE = 15
MAX_CHARS_PER_LINE = 30

MAX_SOURCES_PER_SLIDE = 8

class PresentationBuilder:
    def __init__(self, title: str = "Untitled Presentation"):
        self.title = title
        self.slides = []

    def add_slide(self, slide: SlideData):
        self.slides.append(slide)

    def get_all_unique_sources(self) -> list:
        """Returns a list of all unique Source objects across all slides."""
        seen_ids = set()
        unique_sources = []

        for slide in self.slides:
            for source in slide.sources:
                if source.id not in seen_ids:
                    seen_ids.add(source.id)
                    unique_sources.append(source)

        return unique_sources

    def save_to_file(self, path: str, filename: str):
        prs = Presentation()

        for slide_data in self.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.title

            # Set scaled font size for content
            content_text = slide_data.get_combined_text(MAX_SOURCE_IDS_PER_SLIDE)
            font_size = self._fit_text_to_slide(content_text)
            font_size2 = self._fit_text_to_slide(slide_data.title)

            # Set content text
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = content_text
            p.font.size = font_size

            # Set scaled title font size
            title_shape = slide.shapes.title
            title_paragraph = title_shape.text_frame.paragraphs[0]
            if title_paragraph.runs:
                title_run = title_paragraph.runs[0]
            else:
                title_run = title_paragraph.add_run()
            title_run.font.size = font_size2

        full_path = f"{path}\\{filename}.pptx"
        prs.save(full_path)
        print_(f"Presentation saved to: {full_path}")

    def _fit_text_to_slide(self, text: str) -> Pt:
        """
        Determines the appropriate font size for the slide based on line count.
        Returns the font size in points (Pt).
        """
        line_count = text.count("\n") + 1 + int(len(text) / MAX_CHARS_PER_LINE)

        size = DEFAULT_FONT_SIZE_PT
        while line_count > MAX_LINES_PER_SLIDE and size > MIN_FONT_SIZE_PT:
            size -= FONT_DECREMENT_PT
            line_count = int(line_count * DEFAULT_FONT_SIZE_PT / size)

        if size < MIN_FONT_SIZE_PT:
            raise ValueError("Slide content too long to fit on slide, even at minimum font size.")

        return Pt(size)

    def add_title_slide(self):
        """
        Inserts a title slide at the beginning of the presentation.
        Uses the presentation's title; no content or sources.
        """
        slide = SlideData(
            title= "\n\n" + self.title,
            content="",  # No body text
            sources=[]
        )
        self.slides.insert(0, slide)

    def add_sources_slides(self):
        """Adds one or more slides to the end listing all unique sources."""

        all_sources = self.get_all_unique_sources()
        all_sources.sort(key=lambda s: s.id)  # Sort by ID ascending

        for i in range(0, len(all_sources), MAX_SOURCES_PER_SLIDE):
            chunk = all_sources[i:i + MAX_SOURCES_PER_SLIDE]
            lines = []
            for src in chunk:
                lines.append(f"[{src.id}] {src.title} - {src.link}")

            slide_title = "Sources" if i == 0 else "Sources (Continued)"
            content = "\n".join(lines)

            slide = SlideData(
                title=slide_title,
                content=content,
                sources=[]
            )
            self.slides.append(slide)

    def summarize_contents(self) -> str:
        """
        Returns a formatted string summarizing the internal contents of the presentation.
        Supports both object-based and string-based sources.
        """
        output = [f"Presentation Title: {self.title}", f"Total Slides: {len(self.slides)}", ""]

        for idx, slide in enumerate(self.slides, 1):
            output.append(f"Slide {idx}:")
            output.append(f"  Title: {slide.title}")

            trimmed_content = (slide.content[:100] + '...') if len(slide.content) > 100 else slide.content
            output.append(f"  Content: {trimmed_content or '(empty)'}")

            # Handle sources: either a list of objects or a string
            if isinstance(slide.sources, list) and slide.sources:
                try:
                    source_ids = [f"[{src.id}]" for src in slide.sources]
                    output.append(f"  Sources: {' '.join(source_ids)}")
                except AttributeError:
                    output.append(f"  Sources: [invalid source list structure]")
            elif isinstance(slide.sources, str) and slide.sources.strip():
                output.append(f"  Sources: {slide.sources.strip()}")
            else:
                output.append("  Sources: None")

            output.append("")

        return "\n".join(output)