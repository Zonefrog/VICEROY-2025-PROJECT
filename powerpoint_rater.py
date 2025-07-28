import os
import sys
from pathlib import Path
from pptx import Presentation
from slide_class import SlideData
from presentation_class import PresentationBuilder
import re
from collections import Counter
import statistics
import itertools

# === Config ===
rating_input_mode = 1  # 0 = error, 1 = hardcoded name, 2 = user input, 3 = newest file
configured_filename = "Baking_a_Simple_Cake_Mix_Step-by-Step_Guide.pptx"
results_folder = "results"  # Can be changed

def load_presentation_from_file(path: str) -> PresentationBuilder:
    prs = Presentation(path)
    builder = PresentationBuilder(title="Imported Presentation")

    for slide in prs.slides:
        title = ""
        content = ""
        
        # Get title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # Get content from placeholder 1 if exists (usually body)
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                content = shape.text.strip()
                break

        slide_data = SlideData(
            title=title,
            content=content,
            sources=[]  # No source tracking in raw pptx yet
        )
        main_content, source_ids = extract_slide_content_and_sources(slide_data.content)
        slide_data.content = main_content
        slide_data.sources = source_ids
        builder.add_slide(slide_data)

    builder.slides = [s for s in builder.slides if s.content.strip() != ""]
    builder.slides = [s for s in builder.slides if not is_source_slide(s)]

    return builder

def presentation_to_text_blocks(presentation) -> list[str]:
    """
    Converts slides into blocks of text: 'title\n---\ncontent'
    """
    return [f"{slide.title}\n---\n{slide.content}" for slide in presentation.slides]

def rate_presentation(slide_blocks: list[str]) -> float:
    scores = [
        rate_0(slide_blocks),
        rate_1(slide_blocks),
        rate_2(slide_blocks),
        rate_3(slide_blocks),
        rate_4(slide_blocks)
    ]
    return sum(scores)

def extract_slide_content_and_sources(slide_text: str):
    split_index = slide_text.lower().find("sources")
    if split_index == -1:
        return slide_text.strip(), []

    main_text = slide_text[:split_index].strip()
    sources_text = slide_text[split_index:]

    return main_text, sources_text

def resolve_pptx_path() -> Path:
    # Get working directory
    current_file_path = Path(__file__).resolve()
    current_file_directory = current_file_path.parent
    folder_path = current_file_directory / results_folder

    # Validate folder exists
    if not folder_path.exists() or not folder_path.is_dir():
        raise FileNotFoundError(f"Results folder not found at: {folder_path}")

    # Determine filename based on mode
    if rating_input_mode == 0:
        raise ValueError("Input mode set to error (0). Please reconfigure.")
    
    elif rating_input_mode == 1:
        file_name = configured_filename

    elif rating_input_mode == 2:
        file_name = input("Enter PowerPoint filename (e.g., presentation.pptx): ").strip()

    elif rating_input_mode == 3:
        # Find most recently modified .pptx file
        pptx_files = list(folder_path.glob("*.pptx"))
        if not pptx_files:
            raise FileNotFoundError(f"No .pptx files found in folder: {folder_path}")

        pptx_files.sort(key=lambda f: f.stat().st_mtime, reverse=True)
        file_name = pptx_files[0].name

    else:
        raise ValueError(f"Unknown rating_input_mode: {rating_input_mode}")

    # Combine to full path
    full_path = folder_path / file_name
    print(f"Resolved file path: {full_path}")

    # Validations
    if not full_path.exists():
        raise FileNotFoundError(f"File does not exist: {full_path}")

    if full_path.suffix.lower() != ".pptx":
        raise ValueError(f"File is not a .pptx PowerPoint: {full_path.name}")

    return full_path

def is_source_slide(slide):
    return slide.title.strip().lower() == "sources" or "sources (" in slide.title.lower()

def rate_0(slide_blocks: list[str]) -> float:
    #this is a rating based on good amounts of text, words, and lines.
    #for each of those, we will get a 0->1 rating. Then average, then *2 those to get the final result
    #repeat for every slide.
    target_character_count = 323
    character_forgiveness = 0.15

    target_word_count = 46
    word_forgiveness = 0.15

    target_line_count = 6
    line_forgiveness = 0.35

    scores_list = []

    for slide in slide_blocks:
        character_count = len(slide)
        word_count = len(slide.split(" "))
        line_count = len(slide.split("\n"))

        if character_count >= target_character_count * (1 - character_forgiveness) and character_count <= target_character_count * (1 + character_forgiveness):
            character_score = 1.0
        else:
            character_score = max(0.0, 1.0 - abs(character_count - target_character_count) / (target_character_count * character_forgiveness))

        if word_count >= target_word_count * (1 - word_forgiveness) and word_count <= target_word_count * (1 + word_forgiveness):
            word_score = 1.0
        else:
            word_score = max(0.0, 1.0 - abs(word_count - target_word_count) / (target_word_count * word_forgiveness))

        if line_count >= target_line_count * (1 - line_forgiveness) and line_count <= target_line_count * (1 + line_forgiveness):
            line_score = 1.0
        else:
            line_score = max(0.0, 1.0 - abs(line_count - target_line_count) / (target_line_count * line_forgiveness))

        slide_score = (character_score + word_score + line_score) / 3.0 * 2

        scores_list.append(slide_score)

    return sum(scores_list) / len(scores_list)

def rate_1(slide_blocks: list[str]) -> float:
    """
    Rates the overall presentation based on repetition across all slides.
    Uses n-gram (default bigram) repetition rate over the combined text.
    Returns a score from 0.0 (very repetitive) to 2.0 (no repetition).
    """
    def repetition_score(text, n=2):
        tokens = text.split()
        if len(tokens) < n:
            return 1.0
        ngrams = list(zip(*[tokens[i:] for i in range(n)]))
        total = len(ngrams)
        repeated = sum(1 for count in Counter(ngrams).values() if count > 1)
        return 1.0 - min(1.0, repeated / total) if total > 0 else 1.0

    combined_text = "\n".join(slide_blocks)
    score = repetition_score(combined_text, n=2)

    return score * 2.0

def rate_2(slide_blocks: list[str]) -> float:
    word_counts = [len(slide.split()) for slide in slide_blocks]
    if len(word_counts) <= 1:
        return 1.0  # Not enough data to judge balance

    stdev = statistics.stdev(word_counts)
    avg = statistics.mean(word_counts)
    if avg == 0:
        return 0.0

    # The more balanced, the lower the stdev relative to mean
    balance_ratio = stdev / avg
    score = max(0.0, 1.0 - balance_ratio)  # 1.0 if perfectly balanced

    return score * 2.0  # Scale to 0–2

def rate_3(slide_blocks: list[str]) -> float:
    """
    Evaluates title overlap between slides.
    Extracts titles from slide_blocks (title is before "---").
    Allows up to 50% word overlap before penalizing.
    Returns a score from 0.0 to 2.0.
    """

    def extract_title(block):
        return block.split("---")[0].strip()

    def normalize(title):
        return set(title.lower().split())

    titles = [extract_title(block) for block in slide_blocks]
    norm_titles = [normalize(t) for t in titles if t]

    overlap_penalties = 0
    pair_count = 0

    for t1, t2 in itertools.combinations(norm_titles, 2):
        if not t1 or not t2:
            continue

        overlap = len(t1 & t2)
        min_len = min(len(t1), len(t2))

        if min_len == 0:
            continue

        similarity = overlap / min_len

        if similarity > 0.2:
            overlap_penalties += similarity - 0.2

        pair_count += 1

    if pair_count == 0:
        return 2.0

    avg_penalty = overlap_penalties / pair_count
    score = max(0.0, 1.0 - avg_penalty) * 2.0
    return round(score, 3)

def rate_4(slide_blocks: list[str]) -> float:
    """
    Scores slides based on coherence between title and content.
    Slides with more title words present in the content score higher.
    Returns a 0.0 - 2.0 score.
    """
    def tokenize(text):
        return set(re.findall(r"\b\w+\b", text.lower()))

    scores = []

    for block in slide_blocks:
        parts = block.split("---", 1)
        if len(parts) != 2:
            continue  # Malformed block, skip

        title, content = parts[0].strip(), parts[1].strip()
        title_words = tokenize(title)
        content_words = tokenize(content)

        if not title_words:
            continue  # No basis for comparison

        overlap = title_words & content_words
        overlap_ratio = len(overlap) / len(title_words)

        # Full match = 1.0, no match = 0.0, anything in between scaled
        scores.append(overlap_ratio)

    if not scores:
        return 1.0  # Neutral score if nothing is rated

    return round((sum(scores) / len(scores)) * 2.0, 3)

def run_rating_pipeline():
    try:
        # Step 1: Get file path based on mode
        pptx_path = resolve_pptx_path()
        print(f"Analyzing file: {pptx_path}")

        # Step 2: Load presentation object
        presentation = load_presentation_from_file(pptx_path)
        print("Presentation loaded successfully.\n")

        # Step 3: Convert presentation to slide text blocks
        slide_blocks = presentation_to_text_blocks(presentation)

        # Step 4: Rate it
        total_score = rate_presentation(slide_blocks)
        rating_0 = rate_0(slide_blocks)
        rating_1 = rate_1(slide_blocks)
        rating_2 = rate_2(slide_blocks)
        rating_3 = rate_3(slide_blocks)
        rating_4 = rate_4(slide_blocks)

        # Step 5: Print results
        print("Presentation Summary:\n")
        print(presentation.summarize_contents())

        print("\n--- Rating Results ---")
        print(f"Overall Quality Score: {total_score:.2f} / 10.00")
        print(f"Quality [Amount Correctness] Score: {rating_0:.2f} / 2.00")
        print(f"Quality [Non-Repetition] Score: {rating_1:.2f} / 2.00")
        print(f"Quality [Slide Content Quantity Balance] Score: {rating_2:.2f} / 2.00")
        print(f"Quality [Title Uniquenss] Score: {rating_3:.2f} / 2.00")
        print(f"Quality [Title and Content Matching-ness] Score: {rating_4:.2f} / 2.00")

    except Exception as e:
        print(f"Error during rating process: {str(e)}")

#main caller

if __name__ == "__main__":
    run_rating_pipeline()